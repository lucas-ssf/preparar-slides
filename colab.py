!pip install python-pptx
import pptx
from pptx.util import Inches, Pt

def dividir_texto(texto):
    return [bloco.strip() for bloco in texto.strip().split("\n\n")]


# Texto de entrada
# texto = """[em branco]

# TOMADO PELA MÃO COM JESUS EU VOU
# SIGO-O COMO OVELHA QUE ENCONTROU O PASTOR
# TOMADO PELA MÃO, COM JESUS EU VOU
# AONDE ELE FOR (BIS)

# 1. SE JESUS ME DIZ: AMIGO, DEIXA TUDO E VEM COMIGO
# ONDE TUDO É MAIS FORMOSO E MAIS FELIZ
# SE JESUS ME DIZ: AMIGO, DEIXA TUDO E VEM COMIGO
# EU, MINHA MÃO POREI NA SUA E IREI COM ELE

# TOMADO PELA MÃO COM JESUS EU VOU
# SIGO-O COMO OVELHA QUE ENCONTROU O PASTOR
# TOMADO PELA MÃO, COM JESUS EU VOU
# AONDE ELE FOR (BIS)

# 2. EU TE LEVAREI, AMIGO, A UM LUGAR COMIGO
# ONDE O SOL E AS ESTRELAS BRILHAM MAIS
# EU TE LEVAREI AMIGO A UM LUGAR COMIGO
# ONDE TUDO É MAIS FORMOSO E MAIS FELIZ

# [em branco]
# """

with open("musicas.txt", "r") as file: #<----- Arquivo com as músicas
    texto = file.read()

# Criar a apresentação
prs = Presentation()
blocos = dividir_texto(texto)

for i, bloco in enumerate(blocos, 1):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    new_text_box = slide.shapes.add_textbox(
    left=Inches(0), top=Inches(0), width=Inches(10), height=Inches(8)
    )
    
    for paragraph in new_text_frame.paragraphs:
      paragraph.alignment = PP_ALIGN.CENTER
      for run in paragraph.runs:
        run.font.size = Pt(55) #<------ Tamanho da fonte

    new_text_frame = new_text_box.text_frame
    if bloco=="[em branco]":
      new_text_frame.text = ""
    else:
      new_text_frame.text = bloco
    new_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

pptx_filename = input("Digite o nome dos slides:")
prs.save(pptx_filename+".pptx")

print(f"Arquivo salvo como {pptx_filename}.pptx")
